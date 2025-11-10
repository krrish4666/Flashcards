# flashcards_app.py
from flask import Flask, render_template, request, send_file, redirect, url_for, flash
import os, io, json, requests
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas as pdf_canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

# ------------------------------------------------------------
# Flask setup
# ------------------------------------------------------------
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "local-dev-secret")

# ------------------------------------------------------------
# Load .env variables (ensures key loads locally)
# ------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(BASE_DIR, ".env"))

# ------------------------------------------------------------
# Gemini configuration
# ------------------------------------------------------------
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-pro:generateContent"

print(f"[DEBUG] GEMINI_API_KEY loaded: {bool(GEMINI_KEY)}")

# ------------------------------------------------------------
# Gemini helper (proper REST request for Gemini 2.x)
# ------------------------------------------------------------
def gemini_generate(prompt_text: str) -> str:
    """
    Sends a request to the Gemini 2.5 Pro API endpoint.
    Uses ?key=API_KEY for authentication.
    """
    if not GEMINI_KEY or not GEMINI_URL:
        raise RuntimeError("Gemini API key or URL missing")

    headers = {"Content-Type": "application/json"}

    body = {
        "contents": [
            {
                "role": "user",
                "parts": [{"text": prompt_text}],
            }
        ]
    }

    try:
        resp = requests.post(
            GEMINI_URL,
            headers=headers,
            params={"key": GEMINI_KEY},  # <-- key attached as query param
            json=body,
            timeout=90,
        )

        # Debugging info
        print(f"[DEBUG] Gemini API status: {resp.status_code}")
        if resp.status_code != 200:
            print(f"[DEBUG] Response text: {resp.text}")
        resp.raise_for_status()

        data = resp.json()
        candidates = data.get("candidates", [])
        if not candidates:
            raise ValueError("No candidates returned from Gemini")

        output_text = candidates[0].get("content", {}).get("parts", [{}])[0].get("text", "")

        # Clean up Markdown-style code blocks if present
        if output_text.startswith("```"):
            output_text = output_text.strip("`").strip()
            if output_text.lower().startswith("json"):
                output_text = output_text[4:].strip()

        print("[DEBUG] Gemini raw output:", output_text[:500])
        return output_text

    except Exception as e:
        print(f"[Gemini API error] {e}")
        if hasattr(e, "response") and e.response is not None:
            print("Response content:", e.response.text)
        raise


# ------------------------------------------------------------
# Prompt template
# ------------------------------------------------------------
FLASHCARD_PROMPT = """
You are an intelligent assistant that must create flashcards from the provided content.
Output only valid JSON (no markdown, no code fences, no explanations).

Each flashcard object must have:
  - "question": a short question or heading (string)
  - "answer": 2–4 bullet points or a short summary (string, using \\n for new lines)

Example output:
{{
  "flashcards": [
    {{"question": "What is AI?", "answer": "• Simulation of human intelligence\\n• Enables learning and reasoning"}},
    {{"question": "Applications of AI", "answer": "• Healthcare\\n• Finance\\n• Automation"}}
  ]
}}

Now create flashcards for the content below.
Return only JSON — do not include any text or markdown outside it.

CONTENT:
{input_text}
"""



# ------------------------------------------------------------
# Simple in-memory storage
# ------------------------------------------------------------
STORE = {}

# ------------------------------------------------------------
# Routes
# ------------------------------------------------------------
@app.route("/")
def index():
    sets = [{"id": sid, "title": STORE[sid]["title"]} for sid in STORE]
    return render_template("index.html", sets=sets)


@app.route("/new", methods=["GET", "POST"])
def create():
    if request.method == "POST":
        title = request.form.get("title") or "Untitled"
        text_input = request.form.get("text", "")

        # -------- File upload extraction --------
        uploaded_file = request.files.get("file")
        extracted_text = ""
        if uploaded_file and uploaded_file.filename:
            fname = uploaded_file.filename.lower()
            if fname.endswith(".pdf"):
                import fitz
                with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
                    for page in doc:
                        extracted_text += page.get_text()
            elif fname.endswith(".docx"):
                from docx import Document
                doc = Document(uploaded_file)
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
            elif fname.endswith(".txt"):
                extracted_text = uploaded_file.read().decode(errors="ignore")
            else:
                flash("Unsupported file type. Upload PDF, DOCX, or TXT.", "danger")

        # Combine all text
        full_text = (text_input + "\n" + extracted_text).strip()
        if not full_text:
            flash("Please enter text or upload a document.", "danger")
            return redirect(url_for("create"))

        # -------- Gemini generation --------
        prompt = FLASHCARD_PROMPT.format(input_text=full_text[:8000])

        try:
            model_out = gemini_generate(prompt)
            data = json.loads(model_out)
            cards = data.get("flashcards", [])
            if not cards:
                raise ValueError("Empty flashcard list from Gemini")
        except Exception as e:
            flash(f"LLM error or invalid JSON: {e}", "danger")
            cards = [{"question": f"Sample Q{i+1}", "answer": "Sample A"} for i in range(5)]

        sid = str(len(STORE) + 1)
        STORE[sid] = {"title": title, "cards": cards}
        return redirect(url_for("view_set", set_id=sid))

    return render_template("create.html")


@app.route("/set/<set_id>")
def view_set(set_id):
    s = STORE.get(set_id)
    if not s:
        return "Not found", 404
    return render_template("view.html", set_id=set_id, title=s["title"], cards=s["cards"])


# ------------------------------------------------------------
# Export helpers
# ------------------------------------------------------------
def export_flashcards_pdf(cards, title):
    from textwrap import wrap

    bio = io.BytesIO()
    c = pdf_canvas.Canvas(bio, pagesize=landscape(A4))
    w, h = landscape(A4)

    for i, fc in enumerate(cards, start=1):
        # -------- Question Page --------
        c.setFont("Helvetica-Bold", 24)
        c.drawCentredString(w / 2, h - 100, f"{title} — Card {i}")
        c.setFont("Helvetica", 18)
        text_obj = c.beginText(50, h - 180)
        wrapped_q = wrap(fc["question"], width=100)
        for line in wrapped_q:
            text_obj.textLine(line)
        c.drawText(text_obj)
        c.showPage()

        # -------- Answer Page --------
        c.setFont("Helvetica-Bold", 20)
        c.drawCentredString(w / 2, h - 100, f"{title} — Card {i} (Answer)")
        c.setFont("Helvetica", 14)
        text_obj = c.beginText(70, h - 160)
        lines = fc["answer"].split("\\n")
        for line in lines:
            for wrapped in wrap(line, width=110):
                text_obj.textLine(wrapped)
        c.drawText(text_obj)
        c.showPage()

    c.save()
    bio.seek(0)
    return bio



def export_flashcards_pptx(cards, title):
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()

    for i, fc in enumerate(cards, start=1):
        # -------- Question Slide --------
        slide_q = prs.slides.add_slide(prs.slide_layouts[5])

        # Add question text box
        tx_q = slide_q.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(3))
        tf_q = tx_q.text_frame
        tf_q.word_wrap = True
        p_q = tf_q.paragraphs[0]
        p_q.text = fc["question"]
        p_q.font.size = Pt(32)
        p_q.font.bold = True
        p_q.alignment = PP_ALIGN.CENTER

        # Add footer
        footer_q = slide_q.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        footer_q.text = f"{title} — Card {i} (Question)"

        # -------- Answer Slide --------
        slide_a = prs.slides.add_slide(prs.slide_layouts[5])
        tx_a = slide_a.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(4))
        tf_a = tx_a.text_frame
        tf_a.word_wrap = True

        # Add each bullet line
        for ln in fc["answer"].split("\\n"):
            p = tf_a.add_paragraph()
            p.text = ln.strip()
            p.level = 0
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0, 0, 0)

        # Add footer
        footer_a = slide_a.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        footer_a.text = f"{title} — Card {i} (Answer)"

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio




@app.route("/set/<set_id>/export/pdf")
def export_pdf(set_id):
    s = STORE.get(set_id)
    if not s:
        return "Not found", 404
    bio = export_flashcards_pdf(s["cards"], s["title"])
    return send_file(
        bio,
        mimetype="application/pdf",
        as_attachment=True,
        download_name=f"{s['title']}.pdf",
    )


@app.route("/set/<set_id>/export/pptx")
def export_pptx(set_id):
    s = STORE.get(set_id)
    if not s:
        return "Not found", 404
    bio = export_flashcards_pptx(s["cards"], s["title"])
    return send_file(
        bio,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=f"{s['title']}.pptx",
    )


# ------------------------------------------------------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
